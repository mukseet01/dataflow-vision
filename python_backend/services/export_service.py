
"""
Service for exporting analysis results to various document formats.
"""
import os
import time
import base64
import tempfile
from typing import Dict, Any
import json
from fastapi import BackgroundTasks, HTTPException
import pandas as pd
from datetime import datetime

from models.schemas import ExportRequest, ExportResponse
from utils.temp_files import cleanup_files

async def generate_export(export_request: ExportRequest, background_tasks: BackgroundTasks) -> ExportResponse:
    """
    Generate an export of an analysis in various formats.
    
    Supported formats:
    - pdf: PDF Document
    - pptx: PowerPoint Presentation 
    - xlsx: Excel Spreadsheet
    - docx: Word Document
    """
    start_time = time.time()
    temp_files = []
    
    try:
        # Validate the export format
        if export_request.format not in ['pdf', 'pptx', 'xlsx', 'docx']:
            raise HTTPException(status_code=400, detail="Unsupported export format")
            
        # Get a unique filename for the export
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"{export_request.title.replace(' ', '_')}_{timestamp}.{export_request.format}"
        
        # Create a temporary directory for the export
        output_path = os.path.join(tempfile.gettempdir(), output_filename)
        temp_files.append(output_path)
        
        # Generate the export based on the requested format
        if export_request.format == 'xlsx':
            export_to_excel(export_request.analysis_data, output_path)
        elif export_request.format == 'pdf':
            export_to_pdf(export_request.analysis_data, output_path)
        elif export_request.format == 'pptx':
            export_to_powerpoint(export_request.analysis_data, output_path)
        elif export_request.format == 'docx':
            export_to_word(export_request.analysis_data, output_path)
        
        # In a production environment, you would upload this file to a storage service
        # and return a download URL. For this example, we'll simulate a download URL.
        # In a real-world scenario, this might be a signed URL from S3, Azure, or similar.
        download_url = f"https://example.com/exports/{output_filename}"
        
        # Add cleanup task for temporary files
        background_tasks.add_task(cleanup_files, temp_files)
        
        # Prepare metadata
        metadata = {
            "processing_time": time.time() - start_time,
            "export_timestamp": datetime.now().isoformat(),
            "output_format": export_request.format,
            "file_size": os.path.getsize(output_path) if os.path.exists(output_path) else 0
        }
        
        return ExportResponse(
            request_id=export_request.request_id,
            format=export_request.format,
            export_file_path=output_path,
            download_url=download_url,
            metadata=metadata
        )
        
    except Exception as e:
        # Clean up any temporary files
        background_tasks.add_task(cleanup_files, temp_files)
        raise HTTPException(status_code=500, detail=str(e))

def export_to_excel(analysis_data: Dict[str, Any], output_path: str) -> None:
    """Export analysis data to Excel format."""
    try:
        # Create a Pandas Excel writer
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            # Summary sheet
            summary_data = {
                'Item': ['Analysis Result', 'Processing Time', 'Analysis Date'],
                'Value': [
                    analysis_data.get('analysis_result', 'N/A'),
                    f"{analysis_data.get('metadata', {}).get('processing_time', 0):.2f} seconds",
                    analysis_data.get('metadata', {}).get('processing_timestamp', 'N/A')
                ]
            }
            pd.DataFrame(summary_data).to_excel(writer, sheet_name='Summary', index=False)
            
            # Data Preview sheet if available
            if 'data_summary' in analysis_data and 'head' in analysis_data['data_summary']:
                df = pd.DataFrame(analysis_data['data_summary']['head'])
                df.to_excel(writer, sheet_name='Data Preview', index=False)
                
            # Data Summary sheet
            if 'data_summary' in analysis_data:
                summary = analysis_data['data_summary']
                data = {
                    'Metric': [],
                    'Value': []
                }
                
                if 'shape' in summary:
                    data['Metric'].append('Total Rows')
                    data['Value'].append(summary['shape'][0])
                    data['Metric'].append('Total Columns')
                    data['Value'].append(summary['shape'][1])
                    
                if 'dtypes' in summary:
                    for col, dtype in summary['dtypes'].items():
                        data['Metric'].append(f"Column '{col}' Type")
                        data['Value'].append(dtype)
                        
                if 'missing_values' in summary:
                    for col, count in summary['missing_values'].items():
                        if count > 0:
                            data['Metric'].append(f"Missing in '{col}'")
                            data['Value'].append(count)
                            
                pd.DataFrame(data).to_excel(writer, sheet_name='Data Stats', index=False)
    except Exception as e:
        raise Exception(f"Failed to export to Excel: {str(e)}")

def export_to_pdf(analysis_data: Dict[str, Any], output_path: str) -> None:
    """Export analysis data to PDF format."""
    try:
        import matplotlib.pyplot as plt
        from matplotlib.backends.backend_pdf import PdfPages
        
        with PdfPages(output_path) as pdf:
            # Title page
            plt.figure(figsize=(8.5, 11))
            plt.text(0.5, 0.9, "Data Analysis Report", ha='center', va='center', fontsize=24)
            plt.text(0.5, 0.8, f"Generated on: {analysis_data.get('metadata', {}).get('processing_timestamp', datetime.now().strftime('%Y-%m-%d %H:%M:%S'))}", 
                    ha='center', fontsize=12)
            plt.axis('off')
            pdf.savefig()
            plt.close()
            
            # Analysis result page
            plt.figure(figsize=(8.5, 11))
            plt.text(0.5, 0.95, "Analysis Results", ha='center', fontsize=16)
            
            # Split the analysis text to fit on page
            analysis_text = analysis_data.get('analysis_result', 'No analysis available.')
            text_parts = [analysis_text[i:i+100] for i in range(0, len(analysis_text), 100)]
            
            y_position = 0.85
            for part in text_parts:
                plt.text(0.1, y_position, part, fontsize=10, wrap=True)
                y_position -= 0.05
                if y_position < 0.1:
                    pdf.savefig()
                    plt.close()
                    plt.figure(figsize=(8.5, 11))
                    y_position = 0.9
            
            plt.axis('off')
            pdf.savefig()
            plt.close()
            
            # Data summary page
            if 'data_summary' in analysis_data:
                plt.figure(figsize=(8.5, 11))
                plt.text(0.5, 0.95, "Data Summary", ha='center', fontsize=16)
                
                y_position = 0.85
                summary = analysis_data['data_summary']
                
                if 'shape' in summary:
                    plt.text(0.1, y_position, f"Dataset Shape: {summary['shape'][0]} rows × {summary['shape'][1]} columns", fontsize=10)
                    y_position -= 0.05
                
                if 'columns' in summary:
                    columns_text = f"Columns: {', '.join(summary['columns'][:10])}"
                    if len(summary['columns']) > 10:
                        columns_text += f"... and {len(summary['columns']) - 10} more"
                    plt.text(0.1, y_position, columns_text, fontsize=10)
                    y_position -= 0.05
                
                plt.axis('off')
                pdf.savefig()
                plt.close()
    except Exception as e:
        raise Exception(f"Failed to export to PDF: {str(e)}")

def export_to_powerpoint(analysis_data: Dict[str, Any], output_path: str) -> None:
    """Export analysis data to PowerPoint format."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        # Create a presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Data Analysis Report"
        subtitle.text = f"Generated on {analysis_data.get('metadata', {}).get('processing_timestamp', datetime.now().strftime('%Y-%m-%d %H:%M:%S'))}"
        
        # Analysis results slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = "Analysis Results"
        
        content = slide.placeholders[1]
        text_frame = content.text_frame
        p = text_frame.paragraphs[0]
        p.text = analysis_data.get('analysis_result', 'No analysis available.')
        
        # Data summary slide
        if 'data_summary' in analysis_data:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = "Data Summary"
            
            content = slide.placeholders[1]
            text_frame = content.text_frame
            
            summary = analysis_data['data_summary']
            
            p = text_frame.add_paragraph()
            if 'shape' in summary:
                p.text = f"Dataset Shape: {summary['shape'][0]} rows × {summary['shape'][1]} columns"
            
            p = text_frame.add_paragraph()
            if 'columns' in summary:
                columns_text = f"Columns: {', '.join(summary['columns'][:10])}"
                if len(summary['columns']) > 10:
                    columns_text += f"... and {len(summary['columns']) - 10} more"
                p.text = columns_text
                
            # Add column types
            if 'dtypes' in summary:
                p = text_frame.add_paragraph()
                p.text = "Column Types:"
                
                for col, dtype in list(summary['dtypes'].items())[:5]:  # Show first 5
                    p = text_frame.add_paragraph()
                    p.level = 1
                    p.text = f"{col}: {dtype}"
        
        # Save the presentation
        prs.save(output_path)
    except Exception as e:
        raise Exception(f"Failed to export to PowerPoint: {str(e)}")

def export_to_word(analysis_data: Dict[str, Any], output_path: str) -> None:
    """Export analysis data to Word document format."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        
        # Add title
        title = doc.add_heading('Data Analysis Report', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add date
        date_paragraph = doc.add_paragraph()
        date_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        date_paragraph.add_run(f"Generated on: {analysis_data.get('metadata', {}).get('processing_timestamp', datetime.now().strftime('%Y-%m-%d %H:%M:%S'))}")
        
        # Add a line break
        doc.add_paragraph()
        
        # Analysis results section
        doc.add_heading('Analysis Results', 1)
        doc.add_paragraph(analysis_data.get('analysis_result', 'No analysis available.'))
        
        # Data summary section
        if 'data_summary' in analysis_data:
            doc.add_heading('Data Summary', 1)
            
            summary = analysis_data['data_summary']
            
            if 'shape' in summary:
                doc.add_paragraph(f"Dataset Shape: {summary['shape'][0]} rows × {summary['shape'][1]} columns")
            
            if 'columns' in summary:
                columns_text = f"Columns: {', '.join(summary['columns'][:10])}"
                if len(summary['columns']) > 10:
                    columns_text += f"... and {len(summary['columns']) - 10} more"
                doc.add_paragraph(columns_text)
                
            # Add column types section
            if 'dtypes' in summary:
                doc.add_heading('Column Types', 2)
                table = doc.add_table(rows=1, cols=2)
                table.style = 'Table Grid'
                
                # Add header row
                header_cells = table.rows[0].cells
                header_cells[0].text = 'Column Name'
                header_cells[1].text = 'Data Type'
                
                # Add rows for each column type
                for col, dtype in summary['dtypes'].items():
                    row_cells = table.add_row().cells
                    row_cells[0].text = col
                    row_cells[1].text = str(dtype)
                
            # Add missing values section if there are any
            if 'missing_values' in summary and any(count > 0 for count in summary['missing_values'].values()):
                doc.add_heading('Missing Values', 2)
                table = doc.add_table(rows=1, cols=2)
                table.style = 'Table Grid'
                
                # Add header row
                header_cells = table.rows[0].cells
                header_cells[0].text = 'Column'
                header_cells[1].text = 'Missing Count'
                
                # Add rows for each column with missing values
                for col, count in summary['missing_values'].items():
                    if count > 0:
                        row_cells = table.add_row().cells
                        row_cells[0].text = col
                        row_cells[1].text = str(count)
        
        # Save the document
        doc.save(output_path)
    except Exception as e:
        raise Exception(f"Failed to export to Word: {str(e)}")
